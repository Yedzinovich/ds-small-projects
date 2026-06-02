"""
Generate a 5-minute PowerPoint presentation for the
Intelligent Pharmacy Inventory Management capstone project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG     = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
ACCENT_GREEN= RGBColor(0x28, 0xA7, 0x45)
ACCENT_RED  = RGBColor(0xDC, 0x35, 0x45)
ACCENT_GOLD = RGBColor(0xFF, 0xC1, 0x07)
LIGHT_GRAY  = RGBColor(0xF0, 0xF0, 0xF0)
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)
BLACK       = RGBColor(0x00, 0x00, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ────────────────────────────────────────────────
def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf

def add_rounded_rect(slide, left, top, width, height, fill_color, text,
                     font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = True
    p.font.name = "Calibri"
    shape.text_frame.paragraphs[0].space_before = Pt(4)
    return shape

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 1.5, 1.2, 10, 1.2,
            "Intelligent Pharmacy Inventory Management",
            font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 2.5, 10, 0.8,
            "Minimizing Waste & Preventing Stockouts via\n"
            "Time-Series Forecasting and Cost-Aware Optimization",
            font_size=20, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 4.0, 10, 0.5,
            "Group X  |  Umais Siddiqui  \u2022  David Greer  \u2022  Inna Yedzinovich",
            font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 4.8, 10, 0.5,
            "DATA 622 \u2014 Capstone Project Proposal  |  Spring 2026",
            font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Decorative line
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.6), Inches(5.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem: Optimization Paradox
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 12, 0.8,
            "The Problem: The Optimization Paradox",
            font_size=30, bold=True, color=WHITE)

# Left column: Stockout
add_rounded_rect(slide, 0.8, 1.5, 5.5, 1.2, ACCENT_RED,
                 "\u26a0  STOCKOUTS (Under-stocking)", font_size=18)
add_bullet_list(slide, 0.8, 2.9, 5.5, 2.5, [
    "\u2022  Patient cannot fill prescription \u2192 health risk",
    "\u2022  Lost revenue + emergency reorder cost ($15\u2013$50)",
    "\u2022  Patient churn: switches to competitor pharmacy",
    "\u2022  Regulatory / liability exposure",
], font_size=14, color=WHITE)

# Right column: Wastage
add_rounded_rect(slide, 7.0, 1.5, 5.5, 1.2, ACCENT_GOLD,
                 "\U0001f4b8  WASTAGE (Over-stocking)", font_size=18,
                 font_color=BLACK)
add_bullet_list(slide, 7.0, 2.9, 5.5, 2.5, [
    "\u2022  Medications expire on shelf \u2192 destroyed",
    "\u2022  Holding costs: storage, insurance, refrigeration",
    "\u2022  Capital tied up in unsold inventory",
    "\u2022  Costs independent pharmacies $thousands/year",
], font_size=14, color=WHITE)

# Bottom callout
add_rounded_rect(slide, 2.5, 5.5, 8.3, 1.2, ACCENT_BLUE,
                 "Current methods (Min/Max, Static Averages) ignore seasonality,\n"
                 "weekly patterns, and the asymmetric cost of these two risks.",
                 font_size=15, font_color=WHITE)

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — Literature Review: What Others Are Doing
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 12, 0.8,
            "Literature Review: Pharmacy Inventory Approaches",
            font_size=30, bold=True, color=WHITE)

# Row 1: Traditional
add_rounded_rect(slide, 0.8, 1.4, 3.6, 0.7, MED_GRAY,
                 "Traditional Methods", font_size=16)
add_bullet_list(slide, 0.8, 2.3, 3.6, 2.0, [
    "\u2022  Min/Max reorder points",
    "\u2022  EOQ (Economic Order Qty)",
    "\u2022  ABC/XYZ classification",
    "\u2022  Vendor Managed Inventory",
], font_size=13, color=WHITE, spacing=Pt(4))

# Row 2: ML Approaches
add_rounded_rect(slide, 4.8, 1.4, 3.8, 0.7, ACCENT_BLUE,
                 "ML / Time-Series", font_size=16)
add_bullet_list(slide, 4.8, 2.3, 3.8, 2.0, [
    "\u2022  Rathipriya et al. (2023): Prophet",
    "    for retail pharmacy (12\u201318% MAPE)",
    "\u2022  Oliveira et al. (2023): Prophet vs",
    "    LSTM vs XGBoost for hospitals",
], font_size=13, color=WHITE, spacing=Pt(4))

# Row 3: Industry
add_rounded_rect(slide, 9.0, 1.4, 3.8, 0.7, ACCENT_GREEN,
                 "Industry Practice", font_size=16)
add_bullet_list(slide, 9.0, 2.3, 3.8, 2.0, [
    "\u2022  CVS Health: gradient-boosted",
    "    trees + Rx refill features",
    "\u2022  McKesson/ABC: VMI systems",
    "\u2022  Ghousi et al.: ARIMA for hospitals",
], font_size=13, color=WHITE, spacing=Pt(4))

# Key insight box
add_rounded_rect(slide, 0.8, 4.6, 12.0, 1.0, RGBColor(0x2C, 0x3E, 0x50),
                 "Key Gap: Most approaches optimize for accuracy (MAPE) alone.\n"
                 "Few integrate asymmetric cost functions that reflect real pharmacy economics.",
                 font_size=15, font_color=ACCENT_GOLD)

# Our position
add_textbox(slide, 0.8, 5.9, 12, 1.0,
            "Our Contribution:  Prophet forecasting  +  Asymmetric cost functions  +  "
            "Pharmacy-centric reorder logic",
            font_size=16, bold=True, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — Why MAPE is Not Enough
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 12, 0.8,
            "Why MAPE Is Not Enough",
            font_size=30, bold=True, color=WHITE)

# MAPE formula
add_textbox(slide, 1.0, 1.3, 11, 0.6,
            "MAPE = (1/n) \u2211 |A\u209c \u2013 F\u209c| / A\u209c  \u00d7  100",
            font_size=22, bold=True, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

# The problem
add_textbox(slide, 1.0, 2.1, 11, 0.5,
            "MAPE is symmetric: it penalizes over-prediction and under-prediction equally.",
            font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

# Comparison table via boxes
# Header
add_rounded_rect(slide, 1.0, 2.9, 3.5, 0.6, MED_GRAY,
                 "Error Type", font_size=14)
add_rounded_rect(slide, 4.6, 2.9, 3.5, 0.6, MED_GRAY,
                 "Business Impact", font_size=14)
add_rounded_rect(slide, 8.2, 2.9, 4.0, 0.6, MED_GRAY,
                 "Clinical Impact", font_size=14)

# Row 1: Over-prediction
add_rounded_rect(slide, 1.0, 3.6, 3.5, 0.8, RGBColor(0x2C, 0x3E, 0x50),
                 "Over-predict (+10 units)\n= Wastage", font_size=13, font_color=ACCENT_GOLD)
add_rounded_rect(slide, 4.6, 3.6, 3.5, 0.8, RGBColor(0x2C, 0x3E, 0x50),
                 "Holding cost;\nexpiration risk", font_size=13)
add_rounded_rect(slide, 8.2, 3.6, 4.0, 0.8, RGBColor(0x2C, 0x3E, 0x50),
                 "Minimal", font_size=13, font_color=ACCENT_GREEN)

# Row 2: Under-prediction
add_rounded_rect(slide, 1.0, 4.5, 3.5, 0.8, RGBColor(0x2C, 0x3E, 0x50),
                 "Under-predict (\u221210 units)\n= Stockout", font_size=13, font_color=ACCENT_RED)
add_rounded_rect(slide, 4.6, 4.5, 3.5, 0.8, RGBColor(0x2C, 0x3E, 0x50),
                 "Lost revenue;\nemergency reorder", font_size=13)
add_rounded_rect(slide, 8.2, 4.5, 4.0, 0.8, RGBColor(0x2C, 0x3E, 0x50),
                 "Patient can't fill Rx;\ntherapy interruption", font_size=13, font_color=ACCENT_RED)

# Bottom insight
add_rounded_rect(slide, 2.0, 5.8, 9.3, 0.9, ACCENT_BLUE,
                 "Solution: Define asymmetric cost functions that weight\n"
                 "stockouts 10\u00d7 more heavily than wastage.",
                 font_size=17, font_color=WHITE)


# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — Cost Functions: The Math
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 12, 0.8,
            "Asymmetric Cost Functions",
            font_size=30, bold=True, color=WHITE)

# Wastage function
add_rounded_rect(slide, 0.8, 1.3, 5.8, 0.6, ACCENT_GOLD,
                 "Wastage Cost  C\u1d42", font_size=18, font_color=BLACK)

add_textbox(slide, 0.8, 2.0, 5.8, 0.5,
            "C\u1d42(t) = max(F\u209c \u2013 A\u209c, 0) \u00d7 (c_unit \u00d7 p_exp + c_hold)",
            font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_bullet_list(slide, 0.8, 2.6, 5.8, 2.0, [
    "\u2022  c_unit = drug acquisition cost",
    "\u2022  p_exp = probability of expiration before sale",
    "\u2022  c_hold = daily holding cost per unit",
    "\u2022  Incurred ONLY when forecast > actual demand",
], font_size=13, color=WHITE, spacing=Pt(3))

# Stockout function
add_rounded_rect(slide, 7.0, 1.3, 5.8, 0.6, ACCENT_RED,
                 "Stockout Cost  C\u02e2", font_size=18)

add_textbox(slide, 7.0, 2.0, 5.8, 0.5,
            "C\u02e2(t) = max(A\u209c \u2013 F\u209c, 0) \u00d7 (\u03b1\u00b7c_unit + c_emerg + c_churn)",
            font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_bullet_list(slide, 7.0, 2.6, 5.8, 2.0, [
    "\u2022  \u03b1 = 10 (asymmetric penalty multiplier)",
    "\u2022  c_emerg = emergency reorder cost ($15\u2013$50)",
    "\u2022  c_churn = patient loss revenue impact ($5\u2013$25)",
    "\u2022  Incurred ONLY when actual > forecast demand",
], font_size=13, color=WHITE, spacing=Pt(3))

# Practical example
add_rounded_rect(slide, 0.8, 4.8, 12.0, 0.6, ACCENT_BLUE,
                 "Practical Example: Amoxicillin 500mg  |  10-unit forecast error",
                 font_size=16)

add_rounded_rect(slide, 0.8, 5.6, 5.8, 1.2, RGBColor(0x2C, 0x3E, 0x50),
                 "Over-predict by 10 units:\n"
                 "C\u1d42 = 10 \u00d7 ($0.50 \u00d7 0.05 + $0.03) = $0.55",
                 font_size=14, font_color=ACCENT_GOLD)

add_rounded_rect(slide, 7.0, 5.6, 5.8, 1.2, RGBColor(0x2C, 0x3E, 0x50),
                 "Under-predict by 10 units:\n"
                 "C\u02e2 = 10 \u00d7 (10\u00d7$0.50 + $25 + $10) = $400",
                 font_size=14, font_color=ACCENT_RED)

# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — Technical Architecture
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 12, 0.8,
            "Technical Architecture: Forecast \u2192 Cost-Optimal Reorder",
            font_size=30, bold=True, color=WHITE)

# Pipeline boxes - top row
box_y = 1.6
box_h = 0.9
bw = 2.0
colors_top = [
    (0.8,  "Synthetic\nData Engine",     RGBColor(0xFF, 0xE6, 0xF2)),
    (3.2,  "Data Prep\n& EDA",           RGBColor(0xE6, 0xF2, 0xFF)),
    (5.6,  "Facebook\nProphet",          ACCENT_GREEN),
    (8.0,  "30-Day\nForecast",           RGBColor(0xE6, 0xF2, 0xFF)),
    (10.4, "Reorder\nEngine",            ACCENT_BLUE),
]
for x, txt, clr in colors_top:
    fc = BLACK if clr in [RGBColor(0xFF, 0xE6, 0xF2), RGBColor(0xE6, 0xF2, 0xFF)] else WHITE
    add_rounded_rect(slide, x, box_y, bw, box_h, clr, txt, font_size=13, font_color=fc)

# Arrows - top row
for i in range(len(colors_top) - 1):
    x1 = colors_top[i][0] + bw
    x2 = colors_top[i+1][0]
    mid_y = box_y + box_h / 2
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x1 + 0.05), Inches(mid_y - 0.15), Inches(x2 - x1 - 0.1), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = WHITE
    arrow.line.fill.background()

# Cost layer label
add_textbox(slide, 0.8, 3.0, 12, 0.5,
            "Cost Optimization Layer (NEW)",
            font_size=18, bold=True, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)

# Pipeline boxes - bottom row
box_y2 = 3.5
colors_bot = [
    (3.2,  "Wastage\nCost C\u1d42",     RGBColor(0xFF, 0xF3, 0xCD)),
    (5.6,  "Stockout\nCost C\u02e2",     RGBColor(0xF8, 0xD7, 0xDA)),
    (8.0,  "Total Loss\nL(t)",           RGBColor(0xFF, 0xE6, 0xCC)),
    (10.4, "Cost-Optimal\nOrder Q*",     ACCENT_GREEN),
]
for x, txt, clr in colors_bot:
    fc = BLACK if clr != ACCENT_GREEN else WHITE
    add_rounded_rect(slide, x, box_y2, bw, box_h, clr, txt, font_size=13, font_color=fc)

# Arrows - bottom row
for i in range(len(colors_bot) - 1):
    x1 = colors_bot[i][0] + bw
    x2 = colors_bot[i+1][0]
    mid_y = box_y2 + box_h / 2
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x1 + 0.05), Inches(mid_y - 0.15), Inches(x2 - x1 - 0.1), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MED_GRAY
    arrow.line.fill.background()

# Down arrows from forecast to cost layer
for x_center in [6.6, 9.0]:
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(x_center - 0.15), Inches(2.55), Inches(0.3), Inches(0.8))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT_RED
    arrow.line.fill.background()

# Bottom callout: key insight
add_rounded_rect(slide, 1.5, 4.8, 10.3, 1.0, RGBColor(0x2C, 0x3E, 0x50),
                 "Because \u03b1 = 10, the optimal order Q* lands near the 80th\u201390th percentile\n"
                 "of Prophet's predicted demand distribution \u2192 prioritizes patient safety.",
                 font_size=15, font_color=WHITE)

# Tech stack
add_textbox(slide, 0.8, 6.2, 12, 0.5,
            "Tech Stack:  Python  \u2022  Facebook Prophet  \u2022  pandas  \u2022  Streamlit Dashboard  \u2022  Synthetic Data (HIPAA-safe)",
            font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Who Else Is Doing This? (Time-Series in Pharmacy)
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 12, 0.8,
            "Who Else Is Using Time-Series in Pharmacy?",
            font_size=30, bold=True, color=WHITE)

entries = [
    ("CVS Health (Huang et al., 2021)",
     "Gradient-boosted trees + time-series features for store-level demand.\n"
     "Integrates prescription refill cycle data as exogenous regressors.",
     ACCENT_BLUE),
    ("Retail Pharmacy Chains (Rathipriya et al., 2023)",
     "Facebook Prophet for retail drug demand forecasting.\n"
     "Achieved 12\u201318% MAPE for chronic meds, 25\u201335% for seasonal drugs.",
     ACCENT_GREEN),
    ("Hospital Pharmacies (Oliveira et al., 2023)",
     "Compared Prophet, LSTM, and XGBoost for pharmaceutical consumption.\n"
     "Prophet matched deep learning accuracy with less tuning overhead.",
     RGBColor(0x9B, 0x59, 0xB6)),
    ("Wholesalers: McKesson, AmerisourceBergen",
     "VMI systems using proprietary demand signals and EDI integration.\n"
     "Optimize for distributor logistics, not pharmacy-level clinical priorities.",
     ACCENT_GOLD),
]

y = 1.4
for title, desc, color in entries:
    add_rounded_rect(slide, 0.8, y, 4.0, 0.6, color, title, font_size=13)
    add_textbox(slide, 5.0, y, 7.5, 0.7, desc,
                font_size=13, color=WHITE)
    y += 1.2

# Our differentiator
add_rounded_rect(slide, 1.5, 6.2, 10.3, 0.8, ACCENT_GREEN,
                 "Our Differentiator:  Prophet forecasting  +  Asymmetric cost penalties  "
                 "+  Open-source, pharmacy-centric design",
                 font_size=15, font_color=WHITE)

# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — Evaluation & Timeline
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 12, 0.8,
            "Evaluation Framework & Timeline",
            font_size=30, bold=True, color=WHITE)

# Evaluation table
add_textbox(slide, 0.8, 1.2, 6, 0.5,
            "Dual Evaluation Axes", font_size=20, bold=True, color=ACCENT_GOLD)

eval_items = [
    ("Forecast Accuracy", "MAPE < 20% for chronic drugs; MAE, RMSE", ACCENT_BLUE),
    ("Business Cost", "Minimize asymmetric loss L_total over 30 days", ACCENT_RED),
    ("Service Level", "Fill rate > 95% (% of demand met from stock)", ACCENT_GREEN),
]
y = 1.9
for label, desc, clr in eval_items:
    add_rounded_rect(slide, 0.8, y, 2.5, 0.6, clr, label, font_size=13)
    add_textbox(slide, 3.5, y + 0.1, 4, 0.5, desc, font_size=13, color=WHITE)
    y += 0.75

# Timeline
add_textbox(slide, 7.5, 1.2, 5.5, 0.5,
            "10-Week Timeline", font_size=20, bold=True, color=ACCENT_GOLD)

timeline = [
    ("Wk 1\u20132", "Data generation + EDA"),
    ("Wk 3\u20134", "Prophet training + MAPE eval"),
    ("Wk 5\u20136", "Cost function implementation"),
    ("Wk 7\u20138", "Cost-optimal reorder engine"),
    ("Wk 9\u201310", "Streamlit dashboard + final report"),
]
y = 1.9
for wk, task in timeline:
    add_rounded_rect(slide, 7.5, y, 1.5, 0.5, ACCENT_BLUE, wk, font_size=12)
    add_textbox(slide, 9.2, y + 0.05, 3.8, 0.5, task, font_size=13, color=WHITE)
    y += 0.6

# Deliverables
add_textbox(slide, 0.8, 4.6, 12, 0.5,
            "Key Deliverables", font_size=20, bold=True, color=ACCENT_GOLD)

deliverables = [
    "\u2022  Trained Prophet models for 8 high-volatility medications",
    "\u2022  Asymmetric cost functions (wastage + stockout) integrated into decision engine",
    "\u2022  Interactive Streamlit dashboard with real-time reorder recommendations",
    "\u2022  Comprehensive evaluation: MAPE + asymmetric loss + service-level metrics",
]
add_bullet_list(slide, 0.8, 5.2, 12, 2.0, deliverables,
                font_size=14, color=WHITE, spacing=Pt(4))

# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — Thank You / Q&A
# ════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 1.5, 1.5, 10, 1.0,
            "Thank You",
            font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.8), Inches(4.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_textbox(slide, 1.5, 3.2, 10, 0.7,
            "Questions & Discussion",
            font_size=28, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 4.5, 10, 0.5,
            "Group X  |  Umais Siddiqui  \u2022  David Greer  \u2022  Inna Yedzinovich",
            font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 5.2, 10, 0.5,
            "DATA 622 \u2014 Capstone Project  |  Spring 2026",
            font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Key takeaways
add_bullet_list(slide, 3.0, 5.8, 7.3, 1.2, [
    "Prophet forecasting  \u2192  Asymmetric cost functions  \u2192  Smart reorder decisions",
    "Stockout penalty (\u03b1=10) ensures patient safety is never compromised",
], font_size=13, color=ACCENT_GREEN, spacing=Pt(4))

# ── Save ────────────────────────────────────────────────────────────
output_path = "FinalProjectPresentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
