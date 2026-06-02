"""
Generate a 5-minute MVP Demo PowerPoint presentation for PharmaCast.
Walks through the implemented system and ends with a live dashboard demo.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# -- Colour palette (matches original) --
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG     = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
ACCENT_GREEN= RGBColor(0x28, 0xA7, 0x45)
ACCENT_RED  = RGBColor(0xDC, 0x35, 0x45)
ACCENT_GOLD = RGBColor(0xFF, 0xC1, 0x07)
LIGHT_GRAY  = RGBColor(0xF0, 0xF0, 0xF0)
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_CARD   = RGBColor(0x2C, 0x3E, 0x50)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# -- Helper functions --
def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=WHITE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


def add_rounded_rect(slide, left, top, width, height, fill_color, text,
                     font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = True
    p.font.name = "Calibri"
    shape.text_frame.paragraphs[0].space_before = Pt(4)
    return shape


# ================================================================
# SLIDE 1 -- Title
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 1.5, 1.0, 10, 1.2,
            "PharmaCast: MVP Demo",
            font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 2.4, 10, 0.8,
            "Intelligent Pharmacy Inventory Management\n"
            "Time-Series Forecasting \u2022 Asymmetric Cost Optimization \u2022 Live Dashboard",
            font_size=20, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.5), Inches(5.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_textbox(slide, 1.5, 4.0, 10, 0.5,
            "Group X  |  Umais Siddiqui  \u2022  David Greer  \u2022  Inna Yedzinovich",
            font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 4.8, 10, 0.5,
            "DATA 622 \u2014 Capstone MVP Presentation  |  Spring 2026",
            font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# KPI preview boxes
kpi_y = 5.6
kpis = [
    ("8 Drugs", ACCENT_BLUE),
    ("30-Day Forecast", ACCENT_GREEN),
    ("\u03b1 = 10 Asymmetric Loss", ACCENT_RED),
    ("Live on Cloud Run", ACCENT_GOLD),
]
x = 1.5
for label, clr in kpis:
    fc = BLACK if clr == ACCENT_GOLD else WHITE
    add_rounded_rect(slide, x, kpi_y, 2.5, 0.7, clr, label, font_size=13, font_color=fc)
    x += 2.8


# ================================================================
# SLIDE 2 -- Problem Recap (30 seconds)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 12, 0.8,
            "The Problem We Solved",
            font_size=30, bold=True, color=WHITE)

add_rounded_rect(slide, 0.8, 1.4, 5.5, 1.0, ACCENT_RED,
                 "\u26a0  Stockouts = Patient Harm", font_size=18)
add_bullet_list(slide, 0.8, 2.6, 5.5, 2.0, [
    "\u2022  Patient cannot fill prescription",
    "\u2022  Emergency reorder at premium cost",
    "\u2022  Patient churn to competitor pharmacy",
    "\u2022  Regulatory / liability exposure",
], font_size=14, color=WHITE, spacing=Pt(4))

add_rounded_rect(slide, 7.0, 1.4, 5.5, 1.0, ACCENT_GOLD,
                 "\U0001f4b8  Wastage = Financial Loss", font_size=18, font_color=BLACK)
add_bullet_list(slide, 7.0, 2.6, 5.5, 2.0, [
    "\u2022  Medications expire on shelf",
    "\u2022  Capital tied up in unsold inventory",
    "\u2022  Holding costs: storage, insurance",
    "\u2022  Costs pharmacies $thousands/year",
], font_size=14, color=WHITE, spacing=Pt(4))

add_rounded_rect(slide, 2.0, 5.2, 9.3, 1.2, ACCENT_BLUE,
                 "Our Solution: Prophet forecasting + Asymmetric cost functions\n"
                 "that penalize stockouts 10\u00d7 more than wastage,\n"
                 "deployed as a live Streamlit dashboard on Google Cloud Run.",
                 font_size=16, font_color=WHITE)


# ================================================================
# SLIDE 3 -- What We Built (Architecture)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 12, 0.8,
            "What We Built: End-to-End Architecture",
            font_size=30, bold=True, color=WHITE)

# Pipeline row 1
boxes_top = [
    (0.5,  "Synthetic\nData Engine\n(788 days, 8 drugs)", RGBColor(0xFF, 0xE6, 0xF2), BLACK),
    (3.3,  "Data Prep\n& EDA\n(Validation, Cleaning)", RGBColor(0xE6, 0xF2, 0xFF), BLACK),
    (6.1,  "Facebook\nProphet\n(Per-drug tuning)", ACCENT_GREEN, WHITE),
    (8.9,  "Decision\nEngine\n(Q* Optimization)", ACCENT_BLUE, WHITE),
]
for x, txt, bg, fc in boxes_top:
    add_rounded_rect(slide, x, 1.3, 2.5, 1.4, bg, txt, font_size=12, font_color=fc)

# Arrows top row
for i in range(len(boxes_top) - 1):
    x1 = boxes_top[i][0] + 2.5
    x2 = boxes_top[i+1][0]
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x1 + 0.05), Inches(1.85), Inches(x2 - x1 - 0.1), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = WHITE
    arrow.line.fill.background()

# Pipeline row 2 - deployment
boxes_bot = [
    (0.5,  "Streamlit\nDashboard\n(7 pages)", RGBColor(0xFF, 0x45, 0x00), WHITE),
    (3.3,  "Google OAuth\n+ Admin Panel\n(User Mgmt)", RGBColor(0x9B, 0x59, 0xB6), WHITE),
    (6.1,  "Docker\nContainer\n(python:3.12-slim)", DARK_CARD, WHITE),
    (8.9,  "Google\nCloud Run\n(CI/CD via GitHub)", ACCENT_GOLD, BLACK),
]
for x, txt, bg, fc in boxes_bot:
    add_rounded_rect(slide, x, 3.3, 2.5, 1.4, bg, txt, font_size=12, font_color=fc)

# Down arrow from decision engine to dashboard
arrow = slide.shapes.add_shape(
    MSO_SHAPE.DOWN_ARROW,
    Inches(1.6), Inches(2.75), Inches(0.3), Inches(0.5))
arrow.fill.solid()
arrow.fill.fore_color.rgb = ACCENT_RED
arrow.line.fill.background()

# Arrows bottom row
for i in range(len(boxes_bot) - 1):
    x1 = boxes_bot[i][0] + 2.5
    x2 = boxes_bot[i+1][0]
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x1 + 0.05), Inches(3.85), Inches(x2 - x1 - 0.1), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MED_GRAY
    arrow.line.fill.background()

# Tech stack footer
add_textbox(slide, 0.5, 5.2, 12, 0.5,
            "Tech Stack:  Python 3.12  \u2022  Prophet  \u2022  Streamlit  \u2022  Plotly  \u2022  "
            "Docker  \u2022  Cloud Run  \u2022  GCS  \u2022  GitHub Actions CI/CD",
            font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# File counts
add_rounded_rect(slide, 0.5, 5.9, 3.5, 0.9, DARK_CARD,
                 "11 Python modules\n~3,500 lines of code", font_size=13)
add_rounded_rect(slide, 4.3, 5.9, 3.5, 0.9, DARK_CARD,
                 "Automated pipeline\nData \u2192 Train \u2192 Evaluate \u2192 Deploy", font_size=13)
add_rounded_rect(slide, 8.1, 5.9, 3.5, 0.9, DARK_CARD,
                 "Zero-downtime deploys\nPush to main \u2192 live in ~5 min", font_size=13)


# ================================================================
# SLIDE 4 -- Cost Functions & Q* (The Math)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 12, 0.8,
            "Cost Functions & Optimal Order Quantity Q*",
            font_size=30, bold=True, color=WHITE)

# Wastage
add_rounded_rect(slide, 0.8, 1.3, 5.8, 0.6, ACCENT_GOLD,
                 "Wastage Cost C\u1d42  (Proposal \u00a76.2)", font_size=16, font_color=BLACK)
add_textbox(slide, 0.8, 2.0, 5.8, 0.5,
            "C\u1d42(t) = max(F\u209c \u2013 A\u209c, 0) \u00d7 (c_unit \u00d7 p_exp + c_hold)",
            font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, 0.8, 2.6, 5.8, 1.5, [
    "\u2022  p_exp = 5% (short shelf) / 2% (long shelf)",
    "\u2022  c_hold = $0.03/unit/day",
    "\u2022  Only incurred when forecast > actual",
], font_size=13, color=WHITE, spacing=Pt(3))

# Stockout
add_rounded_rect(slide, 7.0, 1.3, 5.8, 0.6, ACCENT_RED,
                 "Stockout Cost C\u02e2  (Proposal \u00a76.3)", font_size=16)
add_textbox(slide, 7.0, 2.0, 5.8, 0.5,
            "C\u02e2(t) = max(A\u209c \u2013 F\u209c, 0) \u00d7 (\u03b1\u00b7c_unit + c_emerg + c_churn)",
            font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, 7.0, 2.6, 5.8, 1.5, [
    "\u2022  \u03b1 = 10 (asymmetric penalty multiplier)",
    "\u2022  c_emerg = $1.50/unit, c_churn = $5.00/unit",
    "\u2022  Only incurred when actual > forecast",
], font_size=13, color=WHITE, spacing=Pt(3))

# Q* optimization
add_rounded_rect(slide, 0.8, 4.3, 12.0, 0.6, ACCENT_GREEN,
                 "Cost-Optimal Order Quantity Q*  (Newsvendor Model, Proposal \u00a77.2)", font_size=16)
add_textbox(slide, 0.8, 5.0, 12, 0.5,
            "Critical Ratio = C\u02e2 / (C\u02e2 + C\u1d42)     \u2192     Q* = F\u207b\u00b9(CR, \u03bc=forecast_mean, \u03c3=std_est)",
            font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.8, 5.6, 12, 0.8,
            "Because \u03b1=10, the critical ratio pushes Q* to the 80th-90th percentile "
            "of the demand distribution.\nThis means we order MORE than average demand "
            "\u2014 prioritizing patient safety over minimizing holding costs.",
            font_size=14, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)


# ================================================================
# SLIDE 5 -- Model Performance Results
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 12, 0.8,
            "Model Performance: Results",
            font_size=30, bold=True, color=WHITE)

# Results cards row 1
add_rounded_rect(slide, 0.8, 1.4, 3.5, 1.4, ACCENT_BLUE,
                 "Avg wMAPE\n23.6%\n4/8 drugs pass (<20%)", font_size=15)
add_rounded_rect(slide, 4.6, 1.4, 3.5, 1.4, ACCENT_GREEN,
                 "80% CI Coverage\n63.5%\nProbabilistic forecasts", font_size=15)
add_rounded_rect(slide, 8.4, 1.4, 3.8, 1.4, ACCENT_RED,
                 "Asymmetric Loss\nL_total = $17,158\nDrives reorder decisions", font_size=15)

# Drug breakdown
add_textbox(slide, 0.8, 3.2, 12, 0.5,
            "Drug-Level Breakdown", font_size=20, bold=True, color=ACCENT_GOLD)

passing = [
    ("Metformin 500mg", "Diabetes", "Stable", "<20%", ACCENT_GREEN),
    ("Lisinopril 10mg", "Blood Pressure", "Stable", "<20%", ACCENT_GREEN),
    ("Omeprazole 20mg", "GI", "Holiday spike", "<20%", ACCENT_GREEN),
    ("Sertraline 50mg", "Mental Health", "Winter SAD", "<20%", ACCENT_GREEN),
]
failing = [
    ("Amoxicillin 500mg", "Antibiotic", "Winter spike", ">20%", ACCENT_RED),
    ("Albuterol Inhaler", "Respiratory", "Winter spike", ">20%", ACCENT_RED),
    ("Cetirizine 10mg", "Allergy", "Spring spike", ">20%", ACCENT_RED),
    ("Azithromycin 250mg", "Antibiotic", "Winter spike", ">20%", ACCENT_RED),
]

y = 3.9
add_rounded_rect(slide, 0.8, y, 5.5, 0.5, ACCENT_GREEN,
                 "PASSING (<20% wMAPE) \u2014 Chronic/Stable Drugs", font_size=13)
y += 0.6
for name, cat, profile, mape, clr in passing:
    add_textbox(slide, 1.0, y, 5.3, 0.35,
                f"\u2713  {name}  ({cat}, {profile})", font_size=12, color=WHITE)
    y += 0.35

y = 3.9
add_rounded_rect(slide, 7.0, y, 5.5, 0.5, ACCENT_RED,
                 "ABOVE TARGET (>20% wMAPE) \u2014 Seasonal Drugs", font_size=13)
y += 0.6
for name, cat, profile, mape, clr in failing:
    add_textbox(slide, 7.2, y, 5.3, 0.35,
                f"\u2717  {name}  ({cat}, {profile})", font_size=12, color=WHITE)
    y += 0.35

# Literature context
add_rounded_rect(slide, 1.5, 6.0, 10.3, 1.0, DARK_CARD,
                 "Literature Context (Rathipriya et al., 2023):\n"
                 "Chronic drugs: 12\u201318% MAPE  |  Seasonal drugs: 25\u201335% MAPE\n"
                 "Our results are consistent with published benchmarks.",
                 font_size=14, font_color=WHITE)


# ================================================================
# SLIDE 6 -- Authentication & Production Features
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 0.8, 0.3, 12, 0.8,
            "Production-Ready Features",
            font_size=30, bold=True, color=WHITE)

# Feature cards
features = [
    ("Google OAuth\nSign-In", "Secure login with Google\naccount authorization", ACCENT_BLUE, 0.8),
    ("Admin Panel\nUser Management", "Grant access, generate\ninvite codes, manage roles", RGBColor(0x9B, 0x59, 0xB6), 3.6),
    ("Email/Password\nLogin + Invite Codes", "Alternative auth for users\nwithout Google accounts", ACCENT_GREEN, 6.4),
    ("GCS Persistence\nCloud Storage Sync", "Data survives container\nrestarts on Cloud Run", ACCENT_GOLD, 9.2),
]
for title, desc, clr, x in features:
    fc = BLACK if clr == ACCENT_GOLD else WHITE
    add_rounded_rect(slide, x, 1.3, 2.5, 1.2, clr, title, font_size=14, font_color=fc)
    add_textbox(slide, x, 2.7, 2.5, 0.8, desc, font_size=11, color=WHITE, alignment=PP_ALIGN.CENTER)

# Deployment features
add_textbox(slide, 0.8, 3.8, 12, 0.5,
            "Deployment & CI/CD", font_size=20, bold=True, color=ACCENT_GOLD)

deploy_features = [
    "\u2022  GitHub Actions: push to main \u2192 auto-build Docker image \u2192 deploy to Cloud Run",
    "\u2022  Docker container: Python 3.12, pre-trained models baked into image",
    "\u2022  Cloud Run: auto-scaling, 2Gi RAM, 2 vCPU, health checks",
    "\u2022  GCS bucket syncs data/models/outputs for persistence across restarts",
    "\u2022  Upload page: pharmacists can upload their own CSV data and retrain models live",
]
add_bullet_list(slide, 0.8, 4.4, 12, 2.5, deploy_features,
                font_size=14, color=WHITE, spacing=Pt(5))


# ================================================================
# SLIDE 7 -- Live Demo Transition
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 1.5, 1.0, 10, 1.0,
            "Live Demo",
            font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.3), Inches(4.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_textbox(slide, 1.5, 2.8, 10, 0.6,
            "https://pharmacast-ze33vyyu7q-uc.a.run.app",
            font_size=22, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER, bold=True)

# Demo walkthrough plan
add_textbox(slide, 1.5, 3.8, 10, 0.5,
            "Demo Walkthrough", font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

demo_steps = [
    "1.  Google Sign-In \u2192 Authentication flow",
    "2.  Dashboard \u2192 KPIs, reorder recommendations, service level",
    "3.  Forecasts \u2192 30-day Prophet predictions with confidence intervals",
    "4.  Reorder Alerts \u2192 Stock depletion simulation per drug",
    "5.  Cost Analysis \u2192 TCO breakdown, wastage vs understocking scatter",
    "6.  Admin Panel \u2192 User management, invite codes",
]
add_bullet_list(slide, 3.0, 4.4, 7.3, 2.5, demo_steps,
                font_size=16, color=WHITE, spacing=Pt(6))


# ================================================================
# SLIDE 8 -- Thank You / Q&A
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BG)

add_textbox(slide, 1.5, 1.2, 10, 1.0,
            "Thank You",
            font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.6), Inches(4.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_textbox(slide, 1.5, 3.0, 10, 0.7,
            "Questions & Discussion",
            font_size=28, color=ACCENT_GOLD, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 4.2, 10, 0.5,
            "Group X  |  Umais Siddiqui  \u2022  David Greer  \u2022  Inna Yedzinovich",
            font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 5.0, 10, 0.5,
            "DATA 622 \u2014 Capstone MVP  |  Spring 2026",
            font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Key takeaways
add_bullet_list(slide, 2.5, 5.6, 8.3, 1.5, [
    "\u2713  Prophet forecasting with per-drug hyperparameter tuning",
    "\u2713  Asymmetric cost functions (\u03b1=10) prioritize patient safety",
    "\u2713  Full production deployment: Docker + Cloud Run + CI/CD + OAuth",
    "\u2713  Dashboard URL: https://pharmacast-ze33vyyu7q-uc.a.run.app",
], font_size=14, color=ACCENT_GREEN, spacing=Pt(5))


# -- Save --
output_path = "FinalProjectMVP.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
